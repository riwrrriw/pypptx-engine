"""
Shape creation and management for all supported python-pptx shape types
"""
from __future__ import annotations

import os
from typing import Any, Dict, List

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION

from .formatters import FontFormatter, ColorFormatter, LineFormatter, ShadowFormatter
from .flowchart import FlowchartHandler


class ShapeFactory:
    """Factory for creating various types of shapes."""
    
    def __init__(self, color_formatter):
        self.color_formatter = color_formatter
        self.text_handler = TextShapeHandler(color_formatter)
        self.image_handler = ImageShapeHandler()
        self.chart_handler = ChartShapeHandler()
        self.table_handler = TableShapeHandler(color_formatter)
        self.autoshape_handler = AutoShapeHandler(color_formatter)
        self.flowchart_handler = FlowchartHandler(color_formatter)
    
    def create_shape(self, slide, shape_config: Dict[str, Any], base_dir: str) -> None:
        """Create a shape based on configuration."""
        shape_type = shape_config.get("type", "").lower()
        
        # Get position and size
        x = Inches(shape_config.get("x", 0))
        y = Inches(shape_config.get("y", 0))
        w = Inches(shape_config.get("w", 4))
        h = Inches(shape_config.get("h", 1))
        
        if shape_type == "text":
            self.text_handler.create_text_shape(slide, shape_config, x, y, w, h)
        elif shape_type == "bullet":
            self.text_handler.create_bullet_shape(slide, shape_config, x, y, w, h)
        elif shape_type == "image":
            self.image_handler.create_image_shape(slide, shape_config, x, y, w, h, base_dir)
        elif shape_type == "chart":
            self.chart_handler.create_chart_shape(slide, shape_config, x, y, w, h)
        elif shape_type == "table":
            self.table_handler.create_table_shape(slide, shape_config, x, y, w, h)
        elif shape_type == "autoshape":
            self.autoshape_handler.create_autoshape(slide, shape_config, x, y, w, h)
        elif shape_type == "connector":
            self.autoshape_handler.create_connector(slide, shape_config, x, y, w, h)
        elif shape_type == "group":
            self.autoshape_handler.create_group_shape(slide, shape_config, x, y, w, h)
        elif shape_type == "freeform":
            self.autoshape_handler.create_freeform_shape(slide, shape_config, x, y, w, h)
        elif shape_type == "flowchart":
            self.flowchart_handler.create_flowchart(slide, shape_config, x, y, w, h)


class TextShapeHandler:
    """Handle text-based shapes including textboxes and bullet lists."""
    
    def __init__(self, color_formatter):
        self.color_formatter = color_formatter
    
    def create_text_shape(self, slide, config: Dict[str, Any], x, y, w, h) -> None:
        """Create a text box shape with enhanced paragraph and formatting support."""
        textbox = slide.shapes.add_textbox(x, y, w, h)
        text_frame = textbox.text_frame
        
        # Clear default content
        text_frame.clear()
        
        # Set text frame properties
        self._apply_text_frame_formatting(text_frame, config.get("text_frame", {}))
        
        # Add text content with enhanced support
        text_content = config.get("text", "")
        
        if isinstance(text_content, str):
            # Simple text - single paragraph
            self._create_simple_text_paragraph(text_frame, text_content, config)
        elif isinstance(text_content, list):
            # Multiple paragraphs or rich text variants
            self._create_multi_paragraph_text(text_frame, text_content, config)
        elif isinstance(text_content, dict):
            # Rich text with advanced formatting
            self._create_rich_text_content(text_frame, text_content, config)
        
        # Apply shape-level formatting
        self._apply_shape_formatting(textbox, config)
    
    def _create_simple_text_paragraph(self, text_frame, text_content: str, config: Dict[str, Any]) -> None:
        """Create a simple single paragraph with formatting."""
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text_content
        
        # Apply formatting
        FontFormatter.apply_font_formatting(run.font, config.get("font", {}))
        FontFormatter.apply_paragraph_formatting(p, config.get("paragraph", {}))
        
        # Apply hyperlinks and click actions
        self._apply_text_actions(run, config.get("actions", {}))
    
    def _create_multi_paragraph_text(self, text_frame, text_content: List, config: Dict[str, Any]) -> None:
        """Create multiple paragraphs with individual or shared formatting."""
        for i, para_item in enumerate(text_content):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            
            if isinstance(para_item, str):
                # Simple paragraph text
                run = p.add_run()
                run.text = para_item
                
                # Apply shared formatting
                FontFormatter.apply_font_formatting(run.font, config.get("font", {}))
                FontFormatter.apply_paragraph_formatting(p, config.get("paragraph", {}))
                
            elif isinstance(para_item, dict):
                # Paragraph with individual formatting
                para_text = para_item.get("text", "")
                run = p.add_run()
                run.text = para_text
                
                # Apply paragraph-specific formatting first, then fallback to shared
                para_font_config = para_item.get("font", config.get("font", {}))
                para_paragraph_config = para_item.get("paragraph", config.get("paragraph", {}))
                
                FontFormatter.apply_font_formatting(run.font, para_font_config)
                FontFormatter.apply_paragraph_formatting(p, para_paragraph_config)
                
                # Apply paragraph-specific actions
                self._apply_text_actions(run, para_item.get("actions", {}))
    
    def _create_rich_text_content(self, text_frame, text_content: Dict[str, Any], config: Dict[str, Any]) -> None:
        """Create rich text with advanced formatting options."""
        paragraphs_data = text_content.get("paragraphs", [])
        
        for i, para_data in enumerate(paragraphs_data):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            
            # Handle runs within paragraph
            runs_data = para_data.get("runs", [])
            if not runs_data and "text" in para_data:
                # Simple paragraph with just text
                runs_data = [{"text": para_data["text"]}]
            
            for run_data in runs_data:
                run = p.add_run()
                run.text = run_data.get("text", "")
                
                # Apply run-specific formatting
                run_font_config = run_data.get("font", para_data.get("font", config.get("font", {})))
                FontFormatter.apply_font_formatting(run.font, run_font_config)
                
                # Apply run-specific actions
                self._apply_text_actions(run, run_data.get("actions", {}))
            
            # Apply paragraph-level formatting
            para_paragraph_config = para_data.get("paragraph", config.get("paragraph", {}))
            FontFormatter.apply_paragraph_formatting(p, para_paragraph_config)
    
    def _apply_text_actions(self, run, actions_config: Dict[str, Any]) -> None:
        """Apply hyperlinks and click actions to text runs."""
        if not actions_config:
            return
        
        # Hyperlink
        if "hyperlink" in actions_config:
            hyperlink_config = actions_config["hyperlink"]
            url = hyperlink_config.get("url")
            if url:
                run.hyperlink.address = url
        
        # Click action (for shapes with text)
        if "click_action" in actions_config:
            click_config = actions_config["click_action"]
            action_type = click_config.get("type", "").upper()
            
            # Note: Click actions are typically applied to shapes, not text runs
            # This would need to be handled at the shape level
    
    def create_bullet_shape(self, slide, config: Dict[str, Any], x, y, w, h) -> None:
        """Create a bullet list shape."""
        textbox = slide.shapes.add_textbox(x, y, w, h)
        text_frame = textbox.text_frame
        text_frame.clear()
        
        # Set text frame properties
        self._apply_text_frame_formatting(text_frame, config.get("text_frame", {}))
        
        items = config.get("items", [])
        for i, item in enumerate(items):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = str(item)
            p.level = config.get("level", 0)  # Bullet level
            
            # Apply formatting
            FontFormatter.apply_font_formatting(p.font, config.get("font", {}))
            FontFormatter.apply_paragraph_formatting(p, config.get("paragraph", {}))
        
        # Apply shape-level formatting
        self._apply_shape_formatting(textbox, config)
    
    def _apply_text_frame_formatting(self, text_frame, config: Dict[str, Any]) -> None:
        """Apply text frame specific formatting."""
        if not config:
            return
        
        if "margin_left" in config:
            text_frame.margin_left = Inches(config["margin_left"])
        if "margin_right" in config:
            text_frame.margin_right = Inches(config["margin_right"])
        if "margin_top" in config:
            text_frame.margin_top = Inches(config["margin_top"])
        if "margin_bottom" in config:
            text_frame.margin_bottom = Inches(config["margin_bottom"])
        
        if "word_wrap" in config:
            text_frame.word_wrap = bool(config["word_wrap"])
        
        if "auto_size" in config:
            # Handle auto-sizing options
            pass
        
        vertical_anchor = config.get("vertical_anchor", "").lower()
        if vertical_anchor == "top":
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        elif vertical_anchor == "middle":
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        elif vertical_anchor == "bottom":
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM
    
    def _apply_shape_formatting(self, shape, config: Dict[str, Any]) -> None:
        """Apply general shape formatting with transparent text support."""
        # Handle transparent text (no background)
        transparent_text = config.get("transparent", False)
        no_fill = config.get("no_fill", False)
        no_border = config.get("no_border", False)
        
        if transparent_text or no_fill:
            # Remove background fill completely
            shape.fill.background()
        elif "fill" in config:
            # Apply custom fill
            ColorFormatter.apply_fill(shape, config["fill"])
        else:
            # Default: remove fill for clean text appearance
            shape.fill.background()
        
        if transparent_text or no_border:
            # Remove border completely
            shape.line.fill.background()
        elif "line" in config:
            # Apply custom line formatting
            LineFormatter.apply_line_formatting(shape.line, config["line"])
        else:
            # Default: remove border for clean text appearance
            shape.line.fill.background()
        
        # Shadow formatting (still applies even for transparent text)
        if "shadow" in config:
            ShadowFormatter.apply_shadow(shape, config["shadow"])


class ImageShapeHandler:
    """Handle image shapes."""
    
    def create_image_shape(self, slide, config: Dict[str, Any], x, y, w, h, base_dir: str) -> None:
        """Create an image shape."""
        image_path = config.get("path") or config.get("url")
        if not image_path:
            print("[WARN] No image path or URL specified")
            return
        
        # Handle URL or local file path
        temp_image_path = None
        try:
            if image_path.startswith(('http://', 'https://')):
                # Download image from URL
                import requests
                import tempfile
                
                response = requests.get(image_path, stream=True)
                response.raise_for_status()
                
                # Create temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as temp_file:
                    for chunk in response.iter_content(chunk_size=8192):
                        temp_file.write(chunk)
                    temp_image_path = temp_file.name
                    final_image_path = temp_image_path
            else:
                # Resolve local path
                if not os.path.isabs(image_path):
                    image_path = os.path.join(base_dir, image_path)
                
                if not os.path.exists(image_path):
                    print(f"[WARN] Image not found, skipping: {image_path}")
                    return
                
                final_image_path = image_path
            
            # Determine dimensions
            width = Inches(config["w"]) if "w" in config else None
            height = Inches(config["h"]) if "h" in config else None
            
            # Add picture
            picture = slide.shapes.add_picture(final_image_path, x, y, width=width, height=height)
            
            # Apply formatting
            if "line" in config:
                LineFormatter.apply_line_formatting(picture.line, config["line"])
            
            if "shadow" in config:
                ShadowFormatter.apply_shadow(picture, config["shadow"])
            
            # Clean up temporary file if it was downloaded
            if temp_image_path:
                try:
                    os.unlink(temp_image_path)
                except OSError:
                    pass
                    
        except Exception as e:
            print(f"[WARN] Failed to load image: {e}")
            # Create placeholder rectangle instead
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, w, h
            )
            placeholder.text = "Image not available"
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = self.color_formatter.parse_color("#cccccc")


class ChartShapeHandler:
    """Handle chart shapes."""
    
    def create_chart_shape(self, slide, config: Dict[str, Any], x, y, w, h) -> None:
        """Create a chart shape with support for all chart types."""
        chart_type_name = config.get("chartType", "COLUMN_CLUSTERED")
        
        try:
            chart_type = getattr(XL_CHART_TYPE, chart_type_name)
        except AttributeError:
            print(f"[WARN] Unsupported chart type: {chart_type_name}")
            return
        
        # Prepare chart data based on chart type
        chart_data = self._prepare_chart_data(config, chart_type_name)
        if not chart_data:
            return
        
        # Add chart to slide
        chart_shape = slide.shapes.add_chart(chart_type, x, y, w, h, chart_data)
        chart = chart_shape.chart
        
        # Apply chart formatting
        self._apply_chart_formatting(chart, config.get("formatting", {}))
    
    def _prepare_chart_data(self, config: Dict[str, Any], chart_type_name: str):
        """Prepare chart data based on chart type."""
        if chart_type_name in ["XY_SCATTER", "XY_SCATTER_LINES", "XY_SCATTER_LINES_NO_MARKERS", "XY_SCATTER_SMOOTH", "XY_SCATTER_SMOOTH_NO_MARKERS"]:
            # XY/Scatter charts
            chart_data = XyChartData()
            for series_config in config.get("series", []):
                name = series_config.get("name", "Series")
                xy_data = series_config.get("xy_data", [])
                series = chart_data.add_series(name)
                for point in xy_data:
                    if isinstance(point, (list, tuple)) and len(point) >= 2:
                        series.add_data_point(point[0], point[1])
            return chart_data
        
        elif chart_type_name in ["BUBBLE", "BUBBLE_THREE_D_EFFECT"]:
            # Bubble charts
            chart_data = BubbleChartData()
            for series_config in config.get("series", []):
                name = series_config.get("name", "Series")
                bubble_data = series_config.get("bubble_data", [])
                series = chart_data.add_series(name)
                for point in bubble_data:
                    if isinstance(point, (list, tuple)) and len(point) >= 3:
                        series.add_data_point(point[0], point[1], point[2])
            return chart_data
        
        else:
            # Category charts (column, bar, line, pie, area, etc.)
            chart_data = CategoryChartData()
            categories = config.get("categories", [])
            chart_data.categories = categories
            
            for series_config in config.get("series", []):
                name = series_config.get("name", "Series")
                values = series_config.get("values", [])
                chart_data.add_series(name, values)
            return chart_data
    
    def _apply_chart_formatting(self, chart, formatting: Dict[str, Any]) -> None:
        """Apply chart-specific formatting."""
        if not formatting:
            return
        
        # Chart title
        if "title" in formatting:
            chart.has_title = True
            chart.chart_title.text_frame.text = formatting["title"]
        
        # Legend
        if "legend" in formatting:
            legend_config = formatting["legend"]
            chart.has_legend = legend_config.get("visible", True)
            if chart.has_legend:
                legend = chart.legend
                position = legend_config.get("position", "right").upper()
                if hasattr(XL_LEGEND_POSITION, position):
                    legend.position = getattr(XL_LEGEND_POSITION, position)
        
        # Data labels
        if "data_labels" in formatting:
            data_labels_config = formatting["data_labels"]
            for plot in chart.plots:
                plot.has_data_labels = data_labels_config.get("visible", False)
                if plot.has_data_labels:
                    data_labels = plot.data_labels
                    if "position" in data_labels_config:
                        position = data_labels_config["position"].upper()
                        if hasattr(XL_DATA_LABEL_POSITION, position):
                            data_labels.position = getattr(XL_DATA_LABEL_POSITION, position)
        
        # Axes formatting
        if "axes" in formatting:
            axes_config = formatting["axes"]
            if "category" in axes_config and hasattr(chart, 'category_axis'):
                self._format_axis(chart.category_axis, axes_config["category"])
            if "value" in axes_config and hasattr(chart, 'value_axis'):
                self._format_axis(chart.value_axis, axes_config["value"])
    
    def _format_axis(self, axis, axis_config: Dict[str, Any]) -> None:
        """Format chart axis."""
        if "title" in axis_config:
            axis.has_title = True
            axis.axis_title.text_frame.text = axis_config["title"]
        
        if "min_scale" in axis_config:
            axis.minimum_scale = axis_config["min_scale"]
        
        if "max_scale" in axis_config:
            axis.maximum_scale = axis_config["max_scale"]
        
        if "major_unit" in axis_config:
            axis.major_unit = axis_config["major_unit"]
        
        if "minor_unit" in axis_config:
            axis.minor_unit = axis_config["minor_unit"]


class TableShapeHandler:
    """Handle table shapes with comprehensive functionality."""
    
    def __init__(self, color_formatter):
        self.color_formatter = color_formatter
    
    def create_table_shape(self, slide, config: Dict[str, Any], x, y, w, h) -> None:
        """Create a table shape with advanced features."""
        rows = config.get("rows", 2)
        cols = config.get("cols", 2)
        
        # Add table
        table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
        table = table_shape.table
        
        # Fill table data
        data = config.get("data", [])
        self._populate_table_data(table, data, rows, cols)
        
        # Handle merged cells
        merged_cells = config.get("merged_cells", [])
        self._apply_merged_cells(table, merged_cells)
        
        # Apply column widths
        col_widths = config.get("col_widths", [])
        self._apply_column_widths(table, col_widths)
        
        # Apply row heights
        row_heights = config.get("row_heights", [])
        self._apply_row_heights(table, row_heights)
        
        # Apply table-level formatting
        self._apply_table_formatting(table, config.get("formatting", {}))
        
        # Apply cell-specific formatting
        cell_formatting = config.get("cell_formatting", {})
        self._apply_cell_specific_formatting(table, cell_formatting)
        
        # Apply header row formatting
        if config.get("header_row", False):
            self._apply_header_formatting(table, config.get("header_formatting", {}))
    
    def _populate_table_data(self, table, data: List[List], rows: int, cols: int) -> None:
        """Populate table with data."""
        for row_idx, row_data in enumerate(data):
            if row_idx >= rows:
                break
            for col_idx, cell_data in enumerate(row_data):
                if col_idx >= cols:
                    break
                
                cell = table.cell(row_idx, col_idx)
                if isinstance(cell_data, str):
                    cell.text = cell_data
                elif isinstance(cell_data, dict):
                    cell.text = cell_data.get("text", "")
                    # Apply individual cell formatting
                    if "formatting" in cell_data:
                        self._apply_cell_formatting(cell, cell_data["formatting"])
    
    def _apply_merged_cells(self, table, merged_cells: List[Dict[str, Any]]) -> None:
        """Apply cell merging based on configuration."""
        for merge_config in merged_cells:
            start_row = merge_config.get("start_row", 0)
            start_col = merge_config.get("start_col", 0)
            end_row = merge_config.get("end_row", start_row)
            end_col = merge_config.get("end_col", start_col)
            
            try:
                # Get the cells to merge
                start_cell = table.cell(start_row, start_col)
                end_cell = table.cell(end_row, end_col)
                
                # Merge the cells
                start_cell.merge(end_cell)
                
                # Apply merge-specific formatting
                if "formatting" in merge_config:
                    self._apply_cell_formatting(start_cell, merge_config["formatting"])
                    
            except Exception as e:
                print(f"[WARN] Failed to merge cells ({start_row},{start_col}) to ({end_row},{end_col}): {e}")
    
    def _apply_column_widths(self, table, col_widths: List[float]) -> None:
        """Apply column widths."""
        for col_idx, width in enumerate(col_widths):
            if col_idx < len(table.columns):
                table.columns[col_idx].width = Inches(width)
    
    def _apply_row_heights(self, table, row_heights: List[float]) -> None:
        """Apply row heights."""
        for row_idx, height in enumerate(row_heights):
            if row_idx < len(table.rows):
                table.rows[row_idx].height = Inches(height)
    
    def _apply_cell_formatting(self, cell, formatting: Dict[str, Any]) -> None:
        """Apply comprehensive formatting to a table cell."""
        if not formatting:
            return
        
        # Fill formatting
        if "fill" in formatting:
            ColorFormatter.apply_fill(cell, formatting["fill"])
        
        # Border formatting
        if "borders" in formatting:
            self._apply_cell_borders(cell, formatting["borders"])
        
        # Text formatting
        if "font" in formatting:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    FontFormatter.apply_font_formatting(run.font, formatting["font"])
        
        # Text alignment
        if "alignment" in formatting:
            self._apply_text_alignment(cell, formatting["alignment"])
        
        # Margins
        if "margins" in formatting:
            self._apply_cell_margins(cell, formatting["margins"])
    
    def _apply_cell_borders(self, cell, borders: Dict[str, Any]) -> None:
        """Apply border formatting to cell."""
        try:
            # Top border
            if "top" in borders:
                border_config = borders["top"]
                if border_config.get("visible", True):
                    cell.border_top.color = self.color_formatter.parse_color(
                        border_config.get("color", "#000000")
                    )
                    cell.border_top.width = Pt(border_config.get("width", 1))
            
            # Bottom border
            if "bottom" in borders:
                border_config = borders["bottom"]
                if border_config.get("visible", True):
                    cell.border_bottom.color = self.color_formatter.parse_color(
                        border_config.get("color", "#000000")
                    )
                    cell.border_bottom.width = Pt(border_config.get("width", 1))
            
            # Left border
            if "left" in borders:
                border_config = borders["left"]
                if border_config.get("visible", True):
                    cell.border_left.color = self.color_formatter.parse_color(
                        border_config.get("color", "#000000")
                    )
                    cell.border_left.width = Pt(border_config.get("width", 1))
            
            # Right border
            if "right" in borders:
                border_config = borders["right"]
                if border_config.get("visible", True):
                    cell.border_right.color = self.color_formatter.parse_color(
                        border_config.get("color", "#000000")
                    )
                    cell.border_right.width = Pt(border_config.get("width", 1))
                    
        except Exception as e:
            print(f"[WARN] Failed to apply cell borders: {e}")
    
    def _apply_text_alignment(self, cell, alignment: Dict[str, Any]) -> None:
        """Apply text alignment to cell."""
        try:
            from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
            
            for paragraph in cell.text_frame.paragraphs:
                # Horizontal alignment
                h_align = alignment.get("horizontal", "left").upper()
                if hasattr(PP_ALIGN, h_align):
                    paragraph.alignment = getattr(PP_ALIGN, h_align)
            
            # Vertical alignment
            v_align = alignment.get("vertical", "middle").upper()
            if hasattr(MSO_VERTICAL_ANCHOR, v_align):
                cell.vertical_anchor = getattr(MSO_VERTICAL_ANCHOR, v_align)
                
        except Exception as e:
            print(f"[WARN] Failed to apply text alignment: {e}")
    
    def _apply_cell_margins(self, cell, margins: Dict[str, Any]) -> None:
        """Apply margins to cell."""
        try:
            if "left" in margins:
                cell.margin_left = Inches(margins["left"])
            if "right" in margins:
                cell.margin_right = Inches(margins["right"])
            if "top" in margins:
                cell.margin_top = Inches(margins["top"])
            if "bottom" in margins:
                cell.margin_bottom = Inches(margins["bottom"])
        except Exception as e:
            print(f"[WARN] Failed to apply cell margins: {e}")
    
    def _apply_cell_specific_formatting(self, table, cell_formatting: Dict[str, Any]) -> None:
        """Apply formatting to specific cells by coordinates."""
        for cell_key, formatting in cell_formatting.items():
            try:
                # Parse cell coordinates (e.g., "0,1" or "A1")
                if "," in cell_key:
                    row_idx, col_idx = map(int, cell_key.split(","))
                else:
                    # Convert Excel-style reference (A1, B2, etc.)
                    col_idx = ord(cell_key[0].upper()) - ord('A')
                    row_idx = int(cell_key[1:]) - 1
                
                if row_idx < len(table.rows) and col_idx < len(table.columns):
                    cell = table.cell(row_idx, col_idx)
                    self._apply_cell_formatting(cell, formatting)
                    
            except Exception as e:
                print(f"[WARN] Failed to apply formatting to cell {cell_key}: {e}")
    
    def _apply_header_formatting(self, table, header_formatting: Dict[str, Any]) -> None:
        """Apply special formatting to header row."""
        if not header_formatting or len(table.rows) == 0:
            return
        
        # Apply to first row by default
        header_row_idx = header_formatting.get("row", 0)
        if header_row_idx >= len(table.rows):
            return
        
        for col_idx in range(len(table.columns)):
            try:
                cell = table.cell(header_row_idx, col_idx)
                self._apply_cell_formatting(cell, header_formatting)
            except Exception as e:
                print(f"[WARN] Failed to apply header formatting to cell ({header_row_idx},{col_idx}): {e}")
    
    def _apply_table_formatting(self, table, formatting: Dict[str, Any]) -> None:
        """Apply table-level formatting."""
        if not formatting:
            return
        
        # Table style
        if "style" in formatting:
            try:
                # Apply built-in table style if available
                style_name = formatting["style"]
                # Note: python-pptx has limited table style support
                print(f"[INFO] Table style '{style_name}' requested (limited support in python-pptx)")
            except Exception as e:
                print(f"[WARN] Failed to apply table style: {e}")
        
        # Apply formatting to all cells if specified
        if "all_cells" in formatting:
            for row in table.rows:
                for cell in row.cells:
                    self._apply_cell_formatting(cell, formatting["all_cells"])
        
        # Banded rows
        if formatting.get("banded_rows", False):
            self._apply_banded_formatting(table, "rows", formatting.get("band_formatting", {}))
        
        # Banded columns
        if formatting.get("banded_cols", False):
            self._apply_banded_formatting(table, "cols", formatting.get("band_formatting", {}))
    
    def _apply_banded_formatting(self, table, band_type: str, band_formatting: Dict[str, Any]) -> None:
        """Apply alternating row/column formatting."""
        try:
            even_formatting = band_formatting.get("even", {})
            odd_formatting = band_formatting.get("odd", {})
            
            if band_type == "rows":
                for row_idx, row in enumerate(table.rows):
                    formatting = even_formatting if row_idx % 2 == 0 else odd_formatting
                    for cell in row.cells:
                        self._apply_cell_formatting(cell, formatting)
            elif band_type == "cols":
                for col_idx in range(len(table.columns)):
                    formatting = even_formatting if col_idx % 2 == 0 else odd_formatting
                    for row_idx in range(len(table.rows)):
                        cell = table.cell(row_idx, col_idx)
                        self._apply_cell_formatting(cell, formatting)
                        
        except Exception as e:
            print(f"[WARN] Failed to apply banded formatting: {e}")


class AutoShapeHandler:
    """Handle auto shapes and connectors."""
    
    def __init__(self, color_formatter):
        self.color_formatter = color_formatter
    
    def create_autoshape(self, slide, config: Dict[str, Any], x, y, w, h) -> None:
        """Create an auto shape."""
        shape_type_name = config.get("shape_type", "RECTANGLE")
        
        try:
            shape_type = getattr(MSO_SHAPE, shape_type_name)
        except AttributeError:
            print(f"[WARN] Unsupported auto shape type: {shape_type_name}")
            return
        
        # Add auto shape
        autoshape = slide.shapes.add_shape(shape_type, x, y, w, h)
        
        # Add text if specified
        text = config.get("text")
        if text and autoshape.has_text_frame:
            autoshape.text = text
            
            # Apply text formatting
            if "font" in config:
                for paragraph in autoshape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        FontFormatter.apply_font_formatting(run.font, config["font"])
        
        # Apply shape formatting
        if "fill" in config:
            ColorFormatter.apply_fill(autoshape, config["fill"])
        
        if "line" in config:
            LineFormatter.apply_line_formatting(autoshape.line, config["line"])
        
        if "shadow" in config:
            ShadowFormatter.apply_shadow(autoshape, config["shadow"])
    
    def create_connector(self, slide, config: Dict[str, Any], x, y, w, h) -> None:
        """Create a connector shape."""
        connector_type_name = config.get("connector_type", "STRAIGHT")
        
        # Map connector type
        connector_type = MSO_CONNECTOR_TYPE.STRAIGHT
        if hasattr(MSO_CONNECTOR_TYPE, connector_type_name):
            connector_type = getattr(MSO_CONNECTOR_TYPE, connector_type_name)
        
        # Get connection points - x, y, w, h are already Inches() objects
        # Use raw values from config or convert to Inches if specified
        if "begin_x" in config:
            begin_x = Inches(config["begin_x"])
        else:
            begin_x = x
            
        if "begin_y" in config:
            begin_y = Inches(config["begin_y"])
        else:
            begin_y = y
            
        if "end_x" in config:
            end_x = Inches(config["end_x"])
        else:
            end_x = x + w
            
        if "end_y" in config:
            end_y = Inches(config["end_y"])
        else:
            end_y = y + h
        
        # Add connector
        connector = slide.shapes.add_connector(
            connector_type, begin_x, begin_y, end_x, end_y
        )
        
        # Apply line formatting
        if "line" in config:
            LineFormatter.apply_line_formatting(connector.line, config["line"])
    
    def create_group_shape(self, slide, config: Dict[str, Any], x, y, w, h) -> None:
        """Create a group shape containing multiple shapes."""
        # Note: python-pptx doesn't have direct group creation API
        # This is a placeholder for future implementation
        shapes_config = config.get("shapes", [])
        
        # For now, just create individual shapes
        # In a full implementation, you'd need to use lower-level APIs
        for shape_config in shapes_config:
            # Adjust positions relative to group
            shape_x = x + Inches(shape_config.get("x", 0))
            shape_y = y + Inches(shape_config.get("y", 0))
            shape_w = Inches(shape_config.get("w", 1))
            shape_h = Inches(shape_config.get("h", 1))
            
            # Create the shape (this would need the shape factory)
            # For now, just create basic shapes
            if shape_config.get("type") == "autoshape":
                self.create_autoshape(slide, shape_config, shape_x, shape_y, shape_w, shape_h)
    
    def create_freeform_shape(self, slide, config: Dict[str, Any], x, y, w, h) -> None:
        """Create a freeform shape."""
        # Get the freeform builder
        freeform_builder = slide.shapes.build_freeform(x, y)
        
        # Add points from configuration
        points = config.get("points", [])
        for i, point in enumerate(points):
            point_x = x + Inches(point.get("x", 0))
            point_y = y + Inches(point.get("y", 0))
            
            action = point.get("action", "line_to")
            if action == "move_to" or i == 0:
                # First point or explicit move_to
                pass  # Starting point is already set in build_freeform
            elif action == "line_to":
                freeform_builder.add_line_segments([(point_x, point_y)])
            elif action == "curve_to":
                # For curves, need control points
                cp1_x = x + Inches(point.get("cp1_x", 0))
                cp1_y = y + Inches(point.get("cp1_y", 0))
                cp2_x = x + Inches(point.get("cp2_x", 0))
                cp2_y = y + Inches(point.get("cp2_y", 0))
                # Note: python-pptx has limited curve support
                freeform_builder.add_line_segments([(point_x, point_y)])
        
        # Convert to shape
        freeform = freeform_builder.convert_to_shape()
        
        # Apply formatting
        if "fill" in config:
            ColorFormatter.apply_fill(freeform, config["fill"])
        
        if "line" in config:
            LineFormatter.apply_line_formatting(freeform.line, config["line"])
        
        if "shadow" in config:
            ShadowFormatter.apply_shadow(freeform, config["shadow"])
