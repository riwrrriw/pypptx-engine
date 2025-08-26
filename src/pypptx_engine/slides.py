"""
Slide management and layout handling
"""
from __future__ import annotations

import os
import requests
import tempfile
from typing import Any, Dict

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.action import PP_ACTION_TYPE


class SlideManager:
    """Manage slide creation and layout operations."""
    
    def __init__(self, color_formatter):
        self.color_formatter = color_formatter
    
    def create_slide(self, prs: Presentation, slide_config: Dict[str, Any], base_dir: str, shape_factory):
        """Create a slide with specified layout and content."""
        layout_index = slide_config.get("layout", 6)  # Default to blank layout
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        
        # Apply slide background
        self._apply_background(slide, slide_config.get("background"))
        
        # Add shapes to slide
        shapes_config = slide_config.get("shapes", [])
        for shape_config in shapes_config:
            shape_factory.create_shape(slide, shape_config, base_dir)
        
        # Handle placeholders if specified
        placeholders_config = slide_config.get("placeholders", {})
        self._fill_placeholders(slide, placeholders_config)
        
        # Add notes slide if specified
        notes_config = slide_config.get("notes")
        if notes_config:
            self._add_notes_slide(slide, notes_config)
        
        return slide
    
    def _apply_background(self, slide, background_config) -> None:
        """Apply background formatting to slide."""
        if not background_config:
            return
        
        if isinstance(background_config, str):
            # Simple color background
            color = self.color_formatter.parse_color(background_config)
            if color:
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = color
        elif isinstance(background_config, dict):
            fill = slide.background.fill
            bg_type = background_config.get("type", "solid")
            
            if bg_type == "solid":
                fill.solid()
                color = self.color_formatter.parse_color(background_config.get("color"))
                if color:
                    fill.fore_color.rgb = color
            elif bg_type == "gradient":
                fill.gradient()
                # Gradient implementation would go here
            elif bg_type == "picture":
                # Picture background implementation
                image_path = background_config.get("image_path") or background_config.get("url")
                if image_path:
                    self._apply_picture_background(slide, image_path, background_config)
    
    def _fill_placeholders(self, slide, placeholders_config: Dict[str, Any]) -> None:
        """Fill slide placeholders with content."""
        if not placeholders_config:
            return
        
        for placeholder in slide.placeholders:
            placeholder_name = placeholders_config.get(str(placeholder.placeholder_format.idx))
            if placeholder_name and placeholder_name in placeholders_config:
                content = placeholders_config[placeholder_name]
                
                if placeholder.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    if hasattr(placeholder, 'text_frame'):
                        placeholder.text = content.get("text", "")
                    elif hasattr(placeholder, 'insert_picture'):
                        # Handle picture placeholders
                        image_path = content.get("image_path")
                        if image_path:
                            placeholder.insert_picture(image_path)
    
    def _add_notes_slide(self, slide, notes_config: Dict[str, Any]) -> None:
        """Add notes to the slide."""
        if not notes_config:
            return
        
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        
        # Clear existing notes
        notes_text_frame.clear()
        
        # Add notes content
        notes_content = notes_config.get("text", "")
        if isinstance(notes_content, str):
            p = notes_text_frame.paragraphs[0]
            p.text = notes_content
        elif isinstance(notes_content, list):
            for i, para_text in enumerate(notes_content):
                p = notes_text_frame.add_paragraph() if i > 0 else notes_text_frame.paragraphs[0]
                p.text = str(para_text)
        
        # Apply notes formatting if specified
        if "font" in notes_config:
            from .formatters import FontFormatter
            for paragraph in notes_text_frame.paragraphs:
                for run in paragraph.runs:
                    FontFormatter.apply_font_formatting(run.font, notes_config["font"])
    
    def _apply_picture_background(self, slide, image_path: str, config: Dict[str, Any]) -> None:
        """Apply picture background to slide."""
        try:
            # Handle URL or local file path
            if image_path.startswith(('http://', 'https://')):
                # Download image from URL
                response = requests.get(image_path, stream=True)
                response.raise_for_status()
                
                # Create temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as temp_file:
                    for chunk in response.iter_content(chunk_size=8192):
                        temp_file.write(chunk)
                    temp_image_path = temp_file.name
            else:
                # Use local file path
                temp_image_path = image_path
            
            # Add picture as background by creating a full-slide image
            from pptx.util import Inches
            
            # Get actual slide dimensions from presentation
            presentation = slide.part.package.presentation_part.presentation
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height
            
            # Add picture to cover entire slide
            picture = slide.shapes.add_picture(
                temp_image_path, 0, 0, slide_width, slide_height
            )
            
            # Move picture to back (behind all other elements)
            slide.shapes._spTree.remove(picture._element)
            slide.shapes._spTree.insert(2, picture._element)
            
            # Clean up temporary file if it was downloaded
            if image_path.startswith(('http://', 'https://')):
                try:
                    os.unlink(temp_image_path)
                except OSError:
                    pass
                    
        except Exception as e:
            print(f"[WARN] Failed to apply picture background: {e}")
            # Fallback to solid color if specified
            fallback_color = config.get("fallback_color", "#000000")
            color = self.color_formatter.parse_color(fallback_color)
            if color:
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = color
