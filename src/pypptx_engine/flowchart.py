"""
Flowchart creation and management for pypptx-engine
"""
from __future__ import annotations

from typing import Any, Dict, List, Tuple
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT

from .formatters import FontFormatter, ColorFormatter, LineFormatter, ShadowFormatter


class FlowchartHandler:
    """Handle flowchart creation with predefined shapes and automatic connections."""
    
    # Standard flowchart shape mappings
    FLOWCHART_SHAPES = {
        "start": MSO_SHAPE.OVAL,
        "end": MSO_SHAPE.OVAL,
        "process": MSO_SHAPE.RECTANGLE,
        "decision": MSO_SHAPE.DIAMOND,
        "data": MSO_SHAPE.PARALLELOGRAM,
        "document": MSO_SHAPE.FLOWCHART_DOCUMENT,
        "predefined_process": MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS,
        "internal_storage": MSO_SHAPE.FLOWCHART_INTERNAL_STORAGE,
        "manual_input": MSO_SHAPE.FLOWCHART_MANUAL_INPUT,
        "manual_operation": MSO_SHAPE.FLOWCHART_MANUAL_OPERATION,
        "connector": MSO_SHAPE.FLOWCHART_CONNECTOR,
        "off_page_connector": MSO_SHAPE.FLOWCHART_OFFPAGE_CONNECTOR,
        "preparation": MSO_SHAPE.FLOWCHART_PREPARATION,
        "extract": MSO_SHAPE.FLOWCHART_EXTRACT,
        "merge": MSO_SHAPE.FLOWCHART_MERGE,
        "stored_data": MSO_SHAPE.FLOWCHART_STORED_DATA,
        "delay": MSO_SHAPE.FLOWCHART_DELAY,
        "alternate_process": MSO_SHAPE.FLOWCHART_ALTERNATE_PROCESS,
        "card": MSO_SHAPE.FLOWCHART_CARD,
        "punched_tape": MSO_SHAPE.FLOWCHART_PUNCHED_TAPE,
        "summing_junction": MSO_SHAPE.FLOWCHART_SUMMING_JUNCTION,
        "or": MSO_SHAPE.FLOWCHART_OR,
        "collate": MSO_SHAPE.FLOWCHART_COLLATE,
        "sort": MSO_SHAPE.FLOWCHART_SORT,
        "multidocument": MSO_SHAPE.FLOWCHART_MULTIDOCUMENT,
        "terminator": MSO_SHAPE.FLOWCHART_TERMINATOR,
        "display": MSO_SHAPE.FLOWCHART_DISPLAY
    }
    
    # Default styles for different flowchart elements
    DEFAULT_STYLES = {
        "start": {
            "fill": {"type": "solid", "color": "#2ecc71"},
            "line": {"color": "#27ae60", "width": 2},
            "font": {"color": "#ffffff", "bold": True, "size": 14}
        },
        "end": {
            "fill": {"type": "solid", "color": "#e74c3c"},
            "line": {"color": "#c0392b", "width": 2},
            "font": {"color": "#ffffff", "bold": True, "size": 14}
        },
        "process": {
            "fill": {"type": "solid", "color": "#3498db"},
            "line": {"color": "#2980b9", "width": 2},
            "font": {"color": "#ffffff", "bold": True, "size": 12}
        },
        "decision": {
            "fill": {"type": "solid", "color": "#f39c12"},
            "line": {"color": "#e67e22", "width": 2},
            "font": {"color": "#ffffff", "bold": True, "size": 12}
        },
        "data": {
            "fill": {"type": "solid", "color": "#9b59b6"},
            "line": {"color": "#8e44ad", "width": 2},
            "font": {"color": "#ffffff", "bold": True, "size": 12}
        },
        "document": {
            "fill": {"type": "solid", "color": "#1abc9c"},
            "line": {"color": "#16a085", "width": 2},
            "font": {"color": "#ffffff", "bold": True, "size": 12}
        },
        "default": {
            "fill": {"type": "solid", "color": "#95a5a6"},
            "line": {"color": "#7f8c8d", "width": 2},
            "font": {"color": "#2c3e50", "bold": True, "size": 12}
        }
    }
    
    def __init__(self, color_formatter):
        self.color_formatter = color_formatter
        self.created_shapes = {}  # Store created shapes for connection references
    
    def create_flowchart(self, slide, config: Dict[str, Any], x, y, w, h) -> None:
        """Create a complete flowchart from configuration."""
        # Clear previous shapes reference for new flowchart
        self.created_shapes = {}
        
        # Get flowchart elements
        elements = config.get("elements", [])
        connections = config.get("connections", [])
        
        # Create all flowchart elements first
        for element in elements:
            self._create_flowchart_element(slide, element)
        
        # Then create connections between elements
        for connection in connections:
            self._create_connection(slide, connection)
    
    def _create_flowchart_element(self, slide, element_config: Dict[str, Any]) -> None:
        """Create a single flowchart element."""
        element_id = element_config.get("id", "")
        element_type = element_config.get("flowchart_type", "process").lower()
        text = element_config.get("text", "")
        
        # Get position and size
        x = Inches(element_config.get("x", 0))
        y = Inches(element_config.get("y", 0))
        w = Inches(element_config.get("w", 2))
        h = Inches(element_config.get("h", 1))
        
        # Get the appropriate shape type
        shape_type = self.FLOWCHART_SHAPES.get(element_type, MSO_SHAPE.RECTANGLE)
        
        # Create the shape
        shape = slide.shapes.add_shape(shape_type, x, y, w, h)
        
        # Store shape reference for connections
        if element_id:
            self.created_shapes[element_id] = shape
        
        # Add text if specified
        if text and shape.has_text_frame:
            shape.text = text
            text_frame = shape.text_frame
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            text_frame.word_wrap = True
            
            # Center align text
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        # Apply styling
        self._apply_flowchart_styling(shape, element_type, element_config)
    
    def _apply_flowchart_styling(self, shape, element_type: str, config: Dict[str, Any]) -> None:
        """Apply styling to a flowchart element."""
        # Get default style for element type
        default_style = self.DEFAULT_STYLES.get(element_type, self.DEFAULT_STYLES["default"])
        
        # Apply fill formatting
        fill_config = config.get("fill", default_style.get("fill", {}))
        if fill_config:
            ColorFormatter.apply_fill(shape, fill_config)
        
        # Apply line formatting
        line_config = config.get("line", default_style.get("line", {}))
        if line_config:
            LineFormatter.apply_line_formatting(shape.line, line_config)
        
        # Apply font formatting to text
        font_config = config.get("font", default_style.get("font", {}))
        if font_config and shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    FontFormatter.apply_font_formatting(run.font, font_config)
        
        # Apply shadow if specified
        shadow_config = config.get("shadow")
        if shadow_config:
            ShadowFormatter.apply_shadow(shape, shadow_config)
    
    def _create_connection(self, slide, connection_config: Dict[str, Any]) -> None:
        """Create a connection between two flowchart elements."""
        from_id = connection_config.get("from")
        to_id = connection_config.get("to")
        
        if not from_id or not to_id:
            print(f"[WARN] Connection missing from/to IDs: {connection_config}")
            return
        
        from_shape = self.created_shapes.get(from_id)
        to_shape = self.created_shapes.get(to_id)
        
        if not from_shape or not to_shape:
            print(f"[WARN] Could not find shapes for connection: {from_id} -> {to_id}")
            return
        
        # Calculate connection points
        from_point = self._get_connection_point(from_shape, connection_config.get("from_side", "bottom"))
        to_point = self._get_connection_point(to_shape, connection_config.get("to_side", "top"))
        
        # Get connector type
        connector_type_name = connection_config.get("connector_type", "STRAIGHT")
        connector_type = getattr(MSO_CONNECTOR_TYPE, connector_type_name, MSO_CONNECTOR_TYPE.STRAIGHT)
        
        # Create connector
        connector = slide.shapes.add_connector(
            connector_type, 
            from_point[0], from_point[1], 
            to_point[0], to_point[1]
        )
        
        # Apply connector formatting
        line_config = connection_config.get("line", {"color": "#2c3e50", "width": 2})
        LineFormatter.apply_line_formatting(connector.line, line_config)
        
        # Add connection label if specified
        label_text = connection_config.get("label")
        if label_text:
            self._add_connection_label(slide, connector, label_text, connection_config.get("label_config", {}))
    
    def _get_connection_point(self, shape, side: str) -> Tuple[Any, Any]:
        """Get connection point coordinates for a shape side."""
        # Convert Inches objects to raw coordinate values and back to Inches
        from pptx.util import Inches
        
        left = Inches(float(shape.left.inches))
        top = Inches(float(shape.top.inches))
        width = Inches(float(shape.width.inches))
        height = Inches(float(shape.height.inches))
        
        if side == "top":
            return (Inches(left.inches + width.inches / 2), top)
        elif side == "bottom":
            return (Inches(left.inches + width.inches / 2), Inches(top.inches + height.inches))
        elif side == "left":
            return (left, Inches(top.inches + height.inches / 2))
        elif side == "right":
            return (Inches(left.inches + width.inches), Inches(top.inches + height.inches / 2))
        elif side == "top-left":
            return (left, top)
        elif side == "top-right":
            return (Inches(left.inches + width.inches), top)
        elif side == "bottom-left":
            return (left, Inches(top.inches + height.inches))
        elif side == "bottom-right":
            return (Inches(left.inches + width.inches), Inches(top.inches + height.inches))
        else:
            # Default to center
            return (Inches(left.inches + width.inches / 2), Inches(top.inches + height.inches / 2))
    
    def _add_connection_label(self, slide, connector, label_text: str, label_config: Dict[str, Any]) -> None:
        """Add a text label to a connection."""
        # Calculate label position (midpoint of connector)
        begin_x = connector.begin_x
        begin_y = connector.begin_y
        end_x = connector.end_x
        end_y = connector.end_y
        
        mid_x = (begin_x + end_x) / 2
        mid_y = (begin_y + end_y) / 2
        
        # Create text box for label
        label_w = Inches(label_config.get("w", 1))
        label_h = Inches(label_config.get("h", 0.3))
        
        # Adjust position to center the label
        label_x = mid_x - label_w / 2
        label_y = mid_y - label_h / 2
        
        textbox = slide.shapes.add_textbox(label_x, label_y, label_w, label_h)
        text_frame = textbox.text_frame
        text_frame.text = label_text
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # Apply label formatting
        default_font = {"size": 10, "color": "#2c3e50", "bold": True}
        font_config = label_config.get("font", default_font)
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            for run in paragraph.runs:
                FontFormatter.apply_font_formatting(run.font, font_config)
        
        # Add background if specified
        if label_config.get("background"):
            ColorFormatter.apply_fill(textbox, label_config["background"])
        else:
            # Default white background with border
            ColorFormatter.apply_fill(textbox, {"type": "solid", "color": "#ffffff"})
            LineFormatter.apply_line_formatting(textbox.line, {"color": "#bdc3c7", "width": 1})


class FlowchartLayoutManager:
    """Utility class for automatic flowchart layout."""
    
    @staticmethod
    def create_vertical_layout(elements: List[Dict[str, Any]], start_x: float = 2, start_y: float = 1, 
                             spacing_y: float = 1.5, element_width: float = 2.5, element_height: float = 1) -> List[Dict[str, Any]]:
        """Create a vertical flowchart layout."""
        positioned_elements = []
        current_y = start_y
        
        for i, element in enumerate(elements):
            positioned_element = element.copy()
            positioned_element.update({
                "x": start_x,
                "y": current_y,
                "w": element_width,
                "h": element_height
            })
            positioned_elements.append(positioned_element)
            current_y += element_height + spacing_y
        
        return positioned_elements
    
    @staticmethod
    def create_horizontal_layout(elements: List[Dict[str, Any]], start_x: float = 1, start_y: float = 2, 
                               spacing_x: float = 3, element_width: float = 2.5, element_height: float = 1) -> List[Dict[str, Any]]:
        """Create a horizontal flowchart layout."""
        positioned_elements = []
        current_x = start_x
        
        for i, element in enumerate(elements):
            positioned_element = element.copy()
            positioned_element.update({
                "x": current_x,
                "y": start_y,
                "w": element_width,
                "h": element_height
            })
            positioned_elements.append(positioned_element)
            current_x += element_width + spacing_x
        
        return positioned_elements
    
    @staticmethod
    def create_decision_tree_layout(elements: List[Dict[str, Any]], start_x: float = 7, start_y: float = 1,
                                  level_spacing_y: float = 2, branch_spacing_x: float = 4) -> List[Dict[str, Any]]:
        """Create a decision tree layout for flowcharts with branching."""
        positioned_elements = []
        
        # Simple decision tree: assumes first element is root, then branches
        if not elements:
            return positioned_elements
        
        # Position root element
        root = elements[0].copy()
        root.update({
            "x": start_x,
            "y": start_y,
            "w": 2.5,
            "h": 1
        })
        positioned_elements.append(root)
        
        # Position remaining elements in branches
        current_y = start_y + level_spacing_y
        branch_positions = [start_x - branch_spacing_x, start_x + branch_spacing_x]
        
        for i, element in enumerate(elements[1:], 1):
            positioned_element = element.copy()
            branch_index = (i - 1) % 2
            positioned_element.update({
                "x": branch_positions[branch_index],
                "y": current_y + ((i - 1) // 2) * level_spacing_y,
                "w": 2.5,
                "h": 1
            })
            positioned_elements.append(positioned_element)
        
        return positioned_elements
    
    @staticmethod
    def auto_connect_sequential(element_ids: List[str], connector_type: str = "STRAIGHT") -> List[Dict[str, Any]]:
        """Create sequential connections between elements."""
        connections = []
        
        for i in range(len(element_ids) - 1):
            connections.append({
                "from": element_ids[i],
                "to": element_ids[i + 1],
                "connector_type": connector_type,
                "from_side": "bottom",
                "to_side": "top"
            })
        
        return connections
    
    @staticmethod
    def auto_connect_decision_tree(root_id: str, branch_ids: List[str], 
                                 labels: List[str] = None) -> List[Dict[str, Any]]:
        """Create decision tree connections."""
        connections = []
        
        for i, branch_id in enumerate(branch_ids):
            connection = {
                "from": root_id,
                "to": branch_id,
                "connector_type": "STRAIGHT",
                "from_side": "bottom-left" if i % 2 == 0 else "bottom-right",
                "to_side": "top"
            }
            
            if labels and i < len(labels):
                connection["label"] = labels[i]
                connection["label_config"] = {"font": {"size": 9, "color": "#2c3e50"}}
            
            connections.append(connection)
        
        return connections
