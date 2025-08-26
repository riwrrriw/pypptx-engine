"""
Core PPTX Engine - Main orchestrator for JSON to PPTX conversion
"""
from __future__ import annotations

import os
from typing import Any, Dict

from pptx import Presentation
from pptx.util import Inches

from .formatters import ColorFormatter
from .shapes import ShapeFactory
from .slides import SlideManager


class PPTXEngine:
    """Main engine for converting JSON specifications to PPTX presentations."""
    
    def __init__(self):
        self.color_formatter = ColorFormatter()
        self.shape_factory = ShapeFactory(self.color_formatter)
        self.slide_manager = SlideManager(self.color_formatter)
    
    def create_presentation(self, config: Dict[str, Any], base_dir: str = "") -> Presentation:
        """Create a PowerPoint presentation from JSON configuration."""
        pres_config = config.get("presentation", {})
        
        # Create presentation
        prs = Presentation()
        
        # Set presentation properties
        self._apply_presentation_properties(prs, pres_config)
        
        # Set slide size
        self._apply_slide_size(prs, pres_config.get("size", {}))
        
        # Create slides
        for slide_config in pres_config.get("slides", []):
            self.slide_manager.create_slide(prs, slide_config, base_dir, self.shape_factory)
        
        return prs
    
    def _apply_presentation_properties(self, prs: Presentation, config: Dict[str, Any]) -> None:
        """Apply presentation-level properties like title, author, etc."""
        properties = config.get("properties", {})
        core_props = prs.core_properties
        
        if "title" in properties:
            core_props.title = properties["title"]
        if "author" in properties:
            core_props.author = properties["author"]
        if "subject" in properties:
            core_props.subject = properties["subject"]
        if "comments" in properties:
            core_props.comments = properties["comments"]
        if "category" in properties:
            core_props.category = properties["category"]
        if "keywords" in properties:
            core_props.keywords = properties["keywords"]
    
    def _apply_slide_size(self, prs: Presentation, size_config: Dict[str, Any]) -> None:
        """Apply custom slide dimensions."""
        if "width_in" in size_config and "height_in" in size_config:
            prs.slide_width = Inches(size_config["width_in"])
            prs.slide_height = Inches(size_config["height_in"])
        elif "width_cm" in size_config and "height_cm" in size_config:
            # Convert cm to inches (1 inch = 2.54 cm)
            prs.slide_width = Inches(size_config["width_cm"] / 2.54)
            prs.slide_height = Inches(size_config["height_cm"] / 2.54)
    
    def generate_presentation(self, config: Dict[str, Any], output_path: str, base_dir: str = "") -> None:
        """Generate and save a PowerPoint presentation from JSON configuration."""
        try:
            # Create presentation
            prs = self.create_presentation(config, base_dir)
            
            # Ensure output directory exists
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            # Save presentation
            prs.save(output_path)
            print(f"✅ Presentation saved to: {output_path}")
            
        except Exception as e:
            print(f"❌ Error generating presentation: {e}")
            raise
